from django.http import HttpResponse
from django.template import RequestContext, loader
from django.shortcuts import render_to_response, redirect
from pptx.chart.data import ChartData
from visitas_stg.models import *
from visitas_stg.tools import *
from datetime import *
from oauth2_provider.views import ProtectedResourceView
from django.contrib.auth.decorators import login_required, user_passes_test
import os, sys
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from django.core.servers.basehttp import FileWrapper
import mimetypes
from django.http import StreamingHttpResponse
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR


from django.db.models import Sum, IntegerField, Q, Count

try:
    import cStringIO as StringIO
except ImportError:
    import StringIO
from xlsxwriter.workbook import Workbook
from django.core.servers.basehttp import FileWrapper
from django.http import StreamingHttpResponse

# Create your views here.
from BuscarVisitas import BuscarVisitas, BuscaVisita

def get_user_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user
    else:
        return None

def register_by_access_token(request):

    #del request.session['access_token']

    if request.session.get('access_token'):
        token = {
        'access_token': request.session.get('access_token'),
        'token_type': 'Bearer'
    }
        return JsonResponse(token)
    else:
        #user = get_user_for_token('3DVteYz9OIH6gvQDyYX78GOpHKXgPy'
        user = request.user
        return get_access_token(user,request)




def get_array_or_none(the_string):
    if the_string is None:
        return None
    else:
        return map(int, the_string.split(','))

def get_array_or_vacio(the_string):
    if the_string =='':
        return None
    else:
        return map(int, the_string.split(','))


def listar_visitas(request):
    usuario = request.user.userprofile
    dependencias = get_array_or_none(request.GET.get('dependencia'))
    if dependencias is None or len(dependencias) == 0:
        if usuario.rol == 'AD':
            dependencias = None
        else:
            dependencias = [usuario.dependencia.id]

    buscador = BuscarVisitas(ids_dependencia=dependencias,
                             rango_fecha_inicio=request.GET.get('fechaInicio', None),
                             rango_fecha_fin=request.GET.get('fechaFin', None),
                             descripcion=request.GET.get('descripcion', None),
                             problematica=request.GET.get('problematica', None),
                             nombre_medio=request.GET.get('nombreMedio', None),
                             ids_region=get_array_or_none(request.GET.get('region')),
                             ids_entidad=get_array_or_none(request.GET.get('estado')),
                             ids_municipio=get_array_or_none(request.GET.get('municipio')),
                             ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                             ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                             ids_partido=get_array_or_none(request.GET.get('partido')),
                             identificador_unico=request.GET.get('identificador_unico', None),
                             ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                             ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                             ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                             ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                             limite_min=request.GET.get('limiteMin', 0),
                             limite_max=request.GET.get('limiteMax', 100),
                             )

    resultados = buscador.buscar()
    json_ans = {}
    json_ans['visitas'] = []
    for visita in resultados['visitas']:
        json_ans['visitas'].append(visita.to_serializable_dict())
    # Exportar a Excel

    output = StringIO.StringIO()
    book = Workbook(output)
    sheet = book.add_worksheet('Visitas')

    if resultados:
        # Add a bold format to use to highlight cells.
        bold = book.add_format({'bold': True})
        # encabezados
        sheet.write(0, 0, "id Unico", bold)
        sheet.write(0, 1, "Dependencia", bold)
        sheet.write(0, 2, "Fecha", bold)
        sheet.write(0, 3, "Region", bold)
        sheet.write(0, 4, "Entidad", bold)
        sheet.write(0, 5, "Municipio", bold)
        sheet.write(0, 6, "Distrito Electoral", bold)
        sheet.write(0, 7, "Partido Gobernante", bold)
        sheet.write(0, 8, "Funcionario", bold)
        sheet.write(0, 9, "Cargo de Funcionario", bold)
        sheet.write(0, 10, "Tipo de Actividad", bold)
        sheet.write(0, 11, "Descripcion", bold)
        sheet.write(0, 12, "Cargo de Funcionario", bold)

        sheet.write(0, 13, "Tipo de Medio", bold)
        sheet.write(0, 14, "Nombre del Medio", bold)
        sheet.write(0, 15, "Capitalizacion", bold)
        sheet.write(0, 16, "Nombre de Participante Local", bold)
        sheet.write(0, 17, "Cargo de Participante Local", bold)
        sheet.write(0, 18, "Problematica Social", bold)

        i = 1
        for obra in json_ans['visitas']:
            sheet.write(i, 0, obra['identificador_unico'])
            sheet.write(i, 1, obra['dependencia']['nombreDependencia'])
            sheet.write(i, 2, obra['fecha_visita'])
            sheet.write(i, 3, obra['region']['numeroRegion'])
            sheet.write(i, 4, obra['entidad']['nombreEstado'])
            sheet.write(i, 5, obra['municipio']['nombreMunicipio'])
            sheet.write(i, 6, obra['distrito_electoral']['nombre_distrito_electoral'])
            sheet.write(i, 7, obra['partido_gobernante']['nombre_partido_gobernante'])
            sheet.write(i, 8, obra['cargo']['nombre_funcionario'])
            sheet.write(i, 9, obra['cargo']['nombre_cargo'])

            for actividad in obra['actividades']:
                sheet.write(i, 10, actividad['tipo_actividad']['nombre_actividad'])
                sheet.write(i, 11, actividad['descripcion'])
                sheet.write(i, 12, actividad['clasificacion']['nombre_clasificacion'])

                j=i
                for capitalizacion in actividad['capitalizaciones']:
                    sheet.write(j, 13, capitalizacion['medio']['nombre_medio'])
                    sheet.write(j, 14, capitalizacion['nombre_medio'])
                    sheet.write(j, 15, capitalizacion['tipo_capitalizacion']['nombre_tipo_capitalizacion'])
                    j += 1
                k=i
                for participantes in actividad['participantes_locales']:
                    sheet.write(k, 16, participantes['nombre'])
                    sheet.write(k, 17, participantes['cargo'])
                    k += 1
                l=i
                for problematica in actividad['problematicas_sociales']:
                    sheet.write(l, 18, problematica['problematica_social'])
                    l += 1
                if j>k and j>l:
                    i=j
                if k>j and k>l:
                    i=k
                if l>j and l>k:
                    i=l
                if j==k and j==l and j==i:
                    i +=1
                else:
                    i=j



        book.close()
    else:
        sheet.write(0, 0,
                   "Los filtros seleccionados no arrojaron informacion alguna sobre las obras, cambie los filtros para una nueva consulta.")
        book.close()

    response = StreamingHttpResponse(FileWrapper(output),
                                    content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = 'attachment; filename="listado_visitas.xlsx"'
    response['Content-Length'] = output.tell()

    output.seek(0)

    return response


def listar_obras(request):


    buscador = BuscarVisitas(ids_dependencia=get_array_or_none(request.GET.get('dependencia')),
                             rango_fecha_inicio=request.GET.get('fechaInicio', None),
                             rango_fecha_fin=request.GET.get('fechaFin', None),
                             descripcion=request.GET.get('descripcion', None),
                             problematica=request.GET.get('problematica', None),
                             nombre_medio=request.GET.get('nombreMedio', None),
                             ids_region=get_array_or_none(request.GET.get('region')),
                             ids_entidad=get_array_or_none(request.GET.get('estado')),
                             ids_municipio=get_array_or_none(request.GET.get('municipio')),
                             ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                             ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                             ids_partido=get_array_or_none(request.GET.get('partido')),
                             identificador_unico=request.GET.get('identificador_unico', None),
                             ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                             ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                             ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                             ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                             limite_min=request.GET.get('limiteMin', 0),
                             limite_max=request.GET.get('limiteMax', 100),
                             )

    resultados = buscador.buscar()

    # Exportar a Excel

    return HttpResponse(resultados)

def buscar_visitas_web(request):
    buscador = BuscarVisitas(ids_dependencia=get_array_or_none(request.GET.get('dependencia')),
                             rango_fecha_inicio=request.GET.get('fechaInicio', None),
                             rango_fecha_fin=request.GET.get('fechaFin', None),
                             descripcion=request.GET.get('descripcion', None),
                             problematica=request.GET.get('problematica', None),
                             nombre_medio=request.GET.get('nombreMedio', None),
                             ids_region=get_array_or_none(request.GET.get('region')),
                             ids_entidad=get_array_or_none(request.GET.get('estado')),
                             ids_municipio=get_array_or_none(request.GET.get('municipio')),
                             ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                             ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                             ids_partido=get_array_or_none(request.GET.get('partido')),
                             identificador_unico=request.GET.get('identificador_unico', None),
                             ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                             ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                             ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                             ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                             limite_min=request.GET.get('limiteMin', 0),
                             limite_max=request.GET.get('limiteMax', 100),
                             )

    resultados = buscador.buscar()
    print resultados

    return HttpResponse(resultados['visitas_stg'])

def catalogos(request):
    return render_to_response('admin/visitas_stg/catalogos.html', locals(),
                              context_instance=RequestContext(request))

def catalogo_funcionarios(request):
    return render_to_response('admin/visitas_stg/catalogo_funcionarios.html', locals(),
                              context_instance=RequestContext(request))

def catalogo_medios(request):
    return render_to_response('admin/visitas_stg/catalogo_medios.html', locals(),
                              context_instance=RequestContext(request))

def catalogo_tipoActividad(request):
    return render_to_response('admin/visitas_stg/catalogo_tipoActividad.html', locals(),
                              context_instance=RequestContext(request))

def catalogo_capitalizacion(request):
    return render_to_response('admin/visitas_stg/catalogo_capitalizacion.html', locals(),
                              context_instance=RequestContext(request))

def consultas(request):
    return render_to_response('admin/visitas_stg/consultas.html', locals(),
                              context_instance=RequestContext(request))
def movimientos(request):
    return render_to_response('admin/visitas_stg/movimientos.html', locals(),
                                 context_instance=RequestContext(request))

def usuarios(request):
    return render_to_response('admin/visitas_stg/usuarios.html', locals(),
                                 context_instance=RequestContext(request))


def consulta_predifinidos(request):

    if request.user.userprofile.rol == 'AD':
        dependencias = Dependencia.objects.all()
    else:
        dependencias = Dependencia.objects.filter(
            Q(id=request.user.userprofile.dependencia_id)
        )

    template = loader.get_template('admin/visitas_stg/consulta_predefinidos/consulta-predefinidos.html')
    context = RequestContext(request, {
        'Regiones': Region.objects.all(),
        'Estados': Estado.objects.all(),
        'Distritos': DistritoElectoral.objects.all(),
        'Municipios': Municipio.objects.all(),
        'Funcionarios': Cargo.objects.all(),
        'Dependencias': dependencias,
        'TipoActividad': TipoActividad.objects.all(),
        'TipoCapitalizacion': TipoCapitalizacion.objects.all(),
        'Medios': Medio.objects.all(),
        'Clasificaciones': Clasificacion.objects.all(),
        'Partidos': PartidoGobernante.objects.all(),
    })
    return HttpResponse(template.render(context))


def consulta_filtros(request):

    print request.user.userprofile.rol

    if request.user.userprofile.rol == 'AD':
        dependencias = Dependencia.objects.all()
    else:
        dependencias = Dependencia.objects.filter(
            Q(id=request.user.userprofile.dependencia_id)
        )

    template = loader.get_template('admin/visitas_stg/consulta_filtros/consulta-filtros.html')
    context = RequestContext(request, {
        'Regiones': Region.objects.all(),
        'Estados': Estado.objects.all(),
        'Distritos': DistritoElectoral.objects.all(),
        'Municipios': Municipio.objects.all(),
        'Funcionarios': Cargo.objects.all(),
        'Dependencias': dependencias,
        'TipoActividad': TipoActividad.objects.all(),
        'TipoCapitalizacion': TipoCapitalizacion.objects.all(),
        'Medios': Medio.objects.all(),
        'Clasificaciones': Clasificacion.objects.all(),
        'Partidos': PartidoGobernante.objects.all(),
    })
    return HttpResponse(template.render(context))


def fichaTecnica(request):
        #prs = Presentation('visitas_stg/static/ppt/fichaTecnica_sisef.pptx')
        prs = Presentation('/home/sisefenlin/visitas/static/ppt/fichaTecnica_sisef.pptx')
        usuario = request.user.userprofile
        buscador = BuscaVisita(
            identificador_unico=request.GET.get('identificador_unico', None)
        )
        resultados = buscador.busca()

        json_map = {}
        json_map['visitas'] = []
        for visita in resultados['visitas']:
            json_map['visitas'].append(visita.to_serializable_dict())

        #generales
        prs.slides[0].shapes[3].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[3].text = json_map['visitas'][0]['identificador_unico']
        prs.slides[0].shapes[4].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[4].text = json_map['visitas'][0]['dependencia']['nombreDependencia']
        prs.slides[0].shapes[5].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[5].text = json_map['visitas'][0]['fecha_visita']
        prs.slides[0].shapes[6].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[6].text = json_map['visitas'][0]['region']['numeroRegion']
        prs.slides[0].shapes[7].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[7].text = json_map['visitas'][0]['entidad']['nombreEstado']
        prs.slides[0].shapes[8].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[8].text = json_map['visitas'][0]['municipio']['nombreMunicipio']
        prs.slides[0].shapes[9].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[9].text = json_map['visitas'][0]['distrito_electoral']['nombre_distrito_electoral']
        prs.slides[0].shapes[10].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[10].text = json_map['visitas'][0]['cargo']['nombre_funcionario']
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[11].text = json_map['visitas'][0]['partido_gobernante']['nombre_partido_gobernante']
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[12].text = json_map['visitas'][0]['cargo']['nombre_cargo']

        prs.slides[0].shapes[22].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[23].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[24].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[22].text = json_map['visitas'][0]['actividades'][0]['tipo_actividad']['nombre_actividad']
        prs.slides[0].shapes[23].text = json_map['visitas'][0]['actividades'][0]['descripcion']
        prs.slides[0].shapes[24].text = json_map['visitas'][0]['actividades'][0]['clasificacion']['nombre_clasificacion']

        #detalles
        prs.slides[0].shapes[15].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[15].text = ""
        prs.slides[0].shapes[16].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[16].text = ""
        if len(json_map['visitas'][0]['actividades'][0]['participantes_locales'])!=0:
            prs.slides[0].shapes[15].text = json_map['visitas'][0]['actividades'][0]['participantes_locales'][0]['nombre']
            prs.slides[0].shapes[16].text = json_map['visitas'][0]['actividades'][0]['participantes_locales'][0]['cargo']

        prs.slides[0].shapes[17].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[17].text = ""
        if len(json_map['visitas'][0]['actividades'][0]['problematicas_sociales'])!=0:
            prs.slides[0].shapes[17].text = json_map['visitas'][0]['actividades'][0]['problematicas_sociales'][0]['problematica_social']

        prs.slides[0].shapes[18].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[18].text = ""
        prs.slides[0].shapes[20].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[20].text = ""
        prs.slides[0].shapes[19].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[19].text = ""
        prs.slides[0].shapes[21].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[21].text = ""

        if len(json_map['visitas'][0]['actividades'][0]['capitalizaciones'])!=0:
            prs.slides[0].shapes[18].text = json_map['visitas'][0]['actividades'][0]['capitalizaciones'][0]['medio']['nombre_medio']
            prs.slides[0].shapes[20].text = json_map['visitas'][0]['actividades'][0]['capitalizaciones'][0]['nombre_medio']
            prs.slides[0].shapes[19].text = json_map['visitas'][0]['actividades'][0]['capitalizaciones'][0]['tipo_capitalizacion']['nombre_tipo_capitalizacion']
            prs.slides[0].shapes[21].text = str(json_map['visitas'][0]['actividades'][0]['capitalizaciones'][0]['cantidad'])


        #prs.save('visitas_stg/static/ppt/ppt-generados/FichaTecnicaVisitas_' + str(usuario.user.id) + '.pptx')
        #the_file = 'visitas_stg/static/ppt/ppt-generados/FichaTecnicaVisitas_' + str(usuario.user.id) + '.pptx'

        prs.save('/home/sisefenlin/visitas/static/ppt/ppt-generados/FichaTecnicaVisitas_' + str(usuario.user.id) + '.pptx')
        the_file = '/home/sisefenlin/visitas/static/ppt/ppt-generados/FichaTecnicaVisitas_' + str(usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response


def Predefinido_Estado(request):

        #prs = Presentation('visitas_stg/static/ppt/Reporte_Estado_sisef.pptx')
        prs = Presentation('/home/sisefenlin/visitas/static/ppt/Reporte_Estado_sisef.pptx')
        usuario = request.user.userprofile

        print request.user.userprofile.rol

        if request.user.userprofile.rol == 'AD':
            dependencias = Dependencia.objects.all()
        else:
            dependencias = Dependencia.objects.filter(
                Q(id=request.user.userprofile.dependencia_id))

        ans = []

        estado = Estado.objects.get(id=request.GET.get('estado_id'))
        medios = Medio.objects.values('id', 'nombre_medio')
        #dependencias = Dependencia.objects.all()
        clasificaciones = Clasificacion.objects.values('id', 'nombre_clasificacion')


        date = datetime.now()
        json_fecha={}
        if date.day >= 10:
            json_fecha['dia'] = str(date.day)
        else:
            json_fecha['dia'] = "0" + str(date.day)
        if date.month >= 10:
            json_fecha['mes'] = str(date.month)
        else:
            json_fecha['mes'] = "0" + str(date.month)
        json_fecha['ano'] = str(date.year)

        map = {}
        map['estado'] = estado.to_serializable_dict()
        map['estado']['distritos_electorales'] = DistritoElectoral.objects.filter(estado_id=estado.id).count()
        map['estado']['municipios'] = Municipio.objects.filter(estado_id=estado.id).count()

        map['dependencias'] = []

        for dependencia in dependencias:
            dependencia_map = dependencia.to_serializable_dict()
            dependencia_map['funcionarios_federales'] = Cargo.objects.filter(dependencia_id=dependencia.id).count()
            dependencia_map['visitas'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            dependencia_map['municipios'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).distinct().count()
            dependencia_map['actividades'] = Actividad.objects.filter(
                Q(visita__dependencia_id=dependencia.id) & Q(visita__entidad_id=dependencia.id)).count()
            dependencia_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).count()
            dependencia_map['capitalizaciones'] = Capitalizacion.objects.filter(
                Q(actividad__visita__entidad_id=estado.id) & Q(
                    actividad__visita__dependencia_id=dependencia.id)).aggregate(Sum('cantidad'))
            map['dependencias'].append(dependencia_map)



        map['medios'] = []
        for medio in medios:
            medio_map = medio

            medio_map['tipos_capitalizacion'] = []
            tipos_capitalizacion = TipoCapitalizacion.objects.all()
            for tipo_capitalizacion in tipos_capitalizacion:
                tipo_map = tipo_capitalizacion.to_serializable_dict()
                tipo_map['numero'] = Capitalizacion.objects.filter(
                    Q(tipo_capitalizacion_id=tipo_capitalizacion.id) & Q(medio_id=medio['id']) & Q(
                        actividad__visita__entidad_id=estado.id)).count()
                medio_map['tipos_capitalizacion'].append(tipo_map)
            map['medios'].append(medio_map)

        map['clasificaciones'] = []
        for clasificacion in clasificaciones:
            tipo_map = clasificacion
            tipo_map['numero'] = Actividad.objects.filter(
                Q(visita__entidad_id=estado.id) & Q(clasificacion_id=clasificacion['id'])).count()
            map['clasificaciones'].append(tipo_map)

        ans.append(map)

        prs.slides[0].shapes[4].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[5].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[6].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[7].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[8].text_frame.paragraphs[0].font.size = Pt(9)


        prs.slides[0].shapes[4].text= json_fecha['dia'] + "/" + json_fecha['mes'] + "/" + json_fecha['ano']
        prs.slides[0].shapes[5].text= '{0:,}'.format(ans[0]['estado']['distritos_electorales'])
        prs.slides[0].shapes[6].text= ans[0]['estado']['nombreEstado']
        #prs.slides[0].shapes[7].text= '{0:,}'.format(ans[0]['estado']['distritos_electorales'])
        prs.slides[0].shapes[7].text= '{0:,}'.format(ans[0]['estado']['municipios'])
        prs.slides[0].shapes[8].text= ans[0]['estado']['region']['numeroRegion']


        table = prs.slides[0].shapes[0].table
        # write body cellstable.cell(1, 0)
        i=1
        totalColumna=total1=total2=total3=total4=total5=0

        for dato in ans[0]['dependencias']:
            if i>=16: break
            table.cell(1, i).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(2, i).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(3, i).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(4, i).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(5, i).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(6, i).text_frame.paragraphs[0].font.size = Pt(8)
            #fill=table.cell(1, i).fill
            #fill.solid()
            #fill.fore_color.rgb = RGBColor(255,0,0)
            #fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
            #fill.fore_color.brightness = 0.25
            #fill.transparency = 0.25
            #fill.background()

            table.cell(1, i).text = str(dato['funcionarios_federales'])
            table.cell(2, i).text = str(dato['visitas'])
            table.cell(3, i).text = str(dato['actividades'])
            table.cell(4, i).text = str(dato['municipios'])
            table.cell(5, i).text = str(dato['participantes_locales'])
            total1=total1 + dato['funcionarios_federales']
            total2=total2 + dato['visitas']
            total3=total3 + dato['actividades']
            total4=total4 + dato['municipios']
            total5=total5 + dato['participantes_locales']
            totalColumna=dato['funcionarios_federales']+dato['visitas']+dato['actividades']+dato['municipios']+dato['participantes_locales']
            table.cell(6, i).text = str(totalColumna)
            i=i+1

        table.cell(1, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(2, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(3, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(4, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(5, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 16).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(1, 16).text = str(total1)
        table.cell(2, 16).text = str(total2)
        table.cell(3, 16).text = str(total3)
        table.cell(4, 16).text = str(total4)
        table.cell(5, 16).text = str(total5)
        table.cell(6, 16).text = str(total1+total2+total3+total4+total5)

        table = prs.slides[0].shapes[1].table
        i=1
        total1=total2=total3=total4=total5=0
        for dato in ans[0]['medios']:
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,2).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,3).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,4).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,5).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,6).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,1).text = str(dato['tipos_capitalizacion'][0]['numero'])
            table.cell(i,2).text = str(dato['tipos_capitalizacion'][1]['numero'])
            table.cell(i,3).text = str(dato['tipos_capitalizacion'][2]['numero'])
            table.cell(i,4).text = str(dato['tipos_capitalizacion'][3]['numero'])
            table.cell(i,5).text = str(dato['tipos_capitalizacion'][4]['numero'])

            totalColumna=0
            for j in  range(0,5):
                totalColumna += dato['tipos_capitalizacion'][j]['numero']
            table.cell(i,6).text = str(totalColumna)

            total1=total1 + dato['tipos_capitalizacion'][0]['numero']
            total2=total2 + dato['tipos_capitalizacion'][1]['numero']
            total3=total3 + dato['tipos_capitalizacion'][2]['numero']
            total4=total4 + dato['tipos_capitalizacion'][3]['numero']
            total5=total5 + dato['tipos_capitalizacion'][4]['numero']


            i=i+1

        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 2).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 3).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 4).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 5).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 6).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text = str(total1)
        table.cell(6, 2).text = str(total2)
        table.cell(6, 3).text = str(total3)
        table.cell(6, 4).text = str(total4)
        table.cell(6, 5).text = str(total5)
        table.cell(6, 6).text = str(total1+total2+total3+total4+total5)

        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'], reverse=True)

        table = prs.slides[0].shapes[2].table
        i=1
        for dato in mayor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreDependencia'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        menor = map['dependencias']
        menor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'])

        table = prs.slides[0].shapes[3].table
        i=1
        for dato in menor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreDependencia'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        total_clasificaciones=0
        for cantidad in ans[0]['clasificaciones']:
            total_clasificaciones = total_clasificaciones + cantidad['numero']

        #grafica pie
        chart_data = ChartData()
        chart_data.categories = [ans[0]['clasificaciones'][0]['nombre_clasificacion'], ans[0]['clasificaciones'][1]['nombre_clasificacion'], ans[0]['clasificaciones'][2]['nombre_clasificacion']]
        chart_data.add_series('Series 1', (float(ans[0]['clasificaciones'][0]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][1]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][2]['numero'])/float(total_clasificaciones)))

        x, y, cx, cy = Inches(6.69), Inches(4.9), Inches(3), Inches(2.5)

        chart = prs.slides[0].shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        #prs.save('visitas_stg/static/ppt/ppt-generados/Reporte_Estado_sisef_' + str(usuario.user.id) + '.pptx')
        #the_file = 'visitas_stg/static/ppt/ppt-generados/Reporte_Estado_sisef_' + str(usuario.user.id) + '.pptx'

        prs.save('/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Estado_sisef_' + str(usuario.user.id) + '.pptx')
        the_file = '/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Estado_sisef_' + str(usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response


def Predefinido_Dependencia(request):

        #prs = Presentation('visitas_stg/static/ppt/Reporte_Dependencia_sisef.pptx')
        prs = Presentation('/home/sisefenlin/visitas/static/ppt/Reporte_Dependencia_sisef.pptx')
        usuario = request.user.userprofile

        print request.user.userprofile.rol

        if request.user.userprofile.rol == 'AD':
            dependencias = Dependencia.objects.all()
        else:
            dependencias = Dependencia.objects.filter(
                Q(id=request.user.userprofile.dependencia_id))

        ans = []

        date = datetime.now()
        json_fecha={}
        if date.day >= 10:
            json_fecha['dia'] = str(date.day)
        else:
            json_fecha['dia'] = "0" + str(date.day)
        if date.month >= 10:
            json_fecha['mes'] = str(date.month)
        else:
            json_fecha['mes'] = "0" + str(date.month)
        json_fecha['ano'] = str(date.year)

        dependencia = Dependencia.objects.get(id=request.GET.get('dependencia_id'))
        visitas = Visita.objects.filter(dependencia_id=dependencia.id)
        clasificaciones = Clasificacion.objects.values('id', 'nombre_clasificacion')
        medios = Medio.objects.all()

        map = {}

        map['dependencia'] = dependencia.to_serializable_dict()
        map['dependencia']['estados_visitados'] = visitas.values('entidad_id').distinct().count()
        map['dependencia']['municipios_visitados'] = visitas.values('entidad_id').distinct().count()
        map['dependencia']['distritos_electorales_visitados'] = visitas.values(
            'distrito_electoral_id').distinct().count()

        map['estados'] = []
        estados = Estado.objects.all()
        for estado in estados:
            estado_map = estado.to_serializable_dict()
            estado_map['total_visitas_funcionarios_federales'] = Visita.objects.filter(
                Q(cargo__dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            estado_map['total_visitas'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            estado_map['total_actividades'] = Actividad.objects.filter(
                Q(visita__entidad_id=estado.id) & Q(visita__dependencia_id=dependencia.id)).count()
            estado_map['municipios'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).values(
                'municipio_id').distinct().count()
            estado_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).count()
            estado_map['capitalizaciones'] = Capitalizacion.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).aggregate(Sum('cantidad'))
            map['estados'].append(estado_map)
        map['estados'].sort(key=lambda e: e['capitalizaciones'])

        map['medios'] = []
        for medio in medios:
            medio_map = medio.to_serializable_dict()

            medio_map['tipos_capitalizacion'] = []
            tipos_capitalizacion = TipoCapitalizacion.objects.all()
            for tipo_capitalizacion in tipos_capitalizacion:
                tipo_map = tipo_capitalizacion.to_serializable_dict()
                tipo_map['numero'] = Capitalizacion.objects.filter(
                    Q(tipo_capitalizacion_id=tipo_capitalizacion.id) & Q(medio_id=medio.id) & Q(
                        actividad__visita__dependencia_id=dependencia.id)).count()
                medio_map['tipos_capitalizacion'].append(tipo_map)
            map['medios'].append(medio_map)

        map['clasificaciones'] = []
        for clasificacion in clasificaciones:
            tipo_map = clasificacion
            tipo_map['numero'] = Actividad.objects.filter(
                Q(visita__dependencia_id=dependencia.id) & Q(clasificacion_id=clasificacion['id'])).count()
            map['clasificaciones'].append(tipo_map)

        ans.append(map)

        prs.slides[0].shapes[8].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[9].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[10].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.size = Pt(9)


        prs.slides[0].shapes[8].text= json_fecha['dia'] + "/" + json_fecha['mes'] + "/" + json_fecha['ano']
        prs.slides[0].shapes[9].text= '{0:,}'.format(ans[0]['dependencia']['distritos_electorales_visitados'])
        prs.slides[0].shapes[10].text= ans[0]['dependencia']['nombreDependencia']
        #prs.slides[0].shapes[7].text= '{0:,}'.format(ans[0]['estado']['distritos_electorales'])
        prs.slides[0].shapes[11].text= '{0:,}'.format(ans[0]['dependencia']['municipios_visitados'])
        prs.slides[0].shapes[12].text= '{0:,}'.format(ans[0]['dependencia']['estados_visitados'])


        table = prs.slides[0].shapes[0].table
        # write body cellstable.cell(1, 0)
        i=1

        mayor = map['estados']
        mayor.sort(key=lambda x: x['total_visitas_funcionarios_federales'], reverse=True)

        total=0

        for dato in ans[0]['estados']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreEstado'])
            table.cell(i, 1).text = str(dato['total_visitas_funcionarios_federales'])
            total += dato['total_visitas_funcionarios_federales']
            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[1].table
        i=1
        mayor = map['estados']
        mayor.sort(key=lambda x: x['total_visitas'], reverse=True)
        total=0
        for dato in ans[0]['estados']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreEstado'])
            table.cell(i, 1).text = str(dato['total_visitas'])
            total+=dato['total_visitas']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[2].table
        i=1
        mayor = map['estados']
        mayor.sort(key=lambda x: x['total_actividades'], reverse=True)
        total=0
        for dato in ans[0]['estados']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreEstado'])
            table.cell(i, 1).text = str(dato['total_actividades'])
            total += dato['total_actividades']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[3].table
        i=1
        mayor = map['estados']
        mayor.sort(key=lambda x: x['municipios'], reverse=True)
        total=0
        for dato in ans[0]['estados']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreEstado'])
            table.cell(i, 1).text = str(dato['municipios'])
            total += dato['municipios']
            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[4].table
        i=1
        mayor = map['estados']
        mayor.sort(key=lambda x: x['participantes_locales'], reverse=True)
        total=0
        for dato in ans[0]['estados']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreEstado'])
            table.cell(i, 1).text = str(dato['participantes_locales'])
            total += dato['participantes_locales']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[5].table
        i=1
        total1=total2=total3=total4=total5=0
        for dato in ans[0]['medios']:
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,2).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,3).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,4).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,5).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,6).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,1).text = str(dato['tipos_capitalizacion'][0]['numero'])
            table.cell(i,2).text = str(dato['tipos_capitalizacion'][1]['numero'])
            table.cell(i,3).text = str(dato['tipos_capitalizacion'][2]['numero'])
            table.cell(i,4).text = str(dato['tipos_capitalizacion'][3]['numero'])
            table.cell(i,5).text = str(dato['tipos_capitalizacion'][4]['numero'])

            totalColumna=0
            for j in  range(0,5):
                totalColumna += dato['tipos_capitalizacion'][j]['numero']
            table.cell(i,6).text = str(totalColumna)

            total1=total1 + dato['tipos_capitalizacion'][0]['numero']
            total2=total2 + dato['tipos_capitalizacion'][1]['numero']
            total3=total3 + dato['tipos_capitalizacion'][2]['numero']
            total4=total4 + dato['tipos_capitalizacion'][3]['numero']
            total5=total5 + dato['tipos_capitalizacion'][4]['numero']
            i=i+1

        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 2).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 3).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 4).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 5).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 6).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text = str(total1)
        table.cell(6, 2).text = str(total2)
        table.cell(6, 3).text = str(total3)
        table.cell(6, 4).text = str(total4)
        table.cell(6, 5).text = str(total5)
        table.cell(6, 6).text = str(total1+total2+total3+total4+total5)

        mayor = map['estados']
        mayor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'], reverse=True)

        table = prs.slides[0].shapes[6].table
        i=1
        for dato in mayor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreEstado'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        menor = map['estados']
        menor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'])

        table = prs.slides[0].shapes[7].table
        i=1
        for dato in menor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreEstado'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        total_clasificaciones=0
        for cantidad in ans[0]['clasificaciones']:
            total_clasificaciones = total_clasificaciones + cantidad['numero']

        #grafica pie
        chart_data = ChartData()
        chart_data.categories = [ans[0]['clasificaciones'][0]['nombre_clasificacion'], ans[0]['clasificaciones'][1]['nombre_clasificacion'], ans[0]['clasificaciones'][2]['nombre_clasificacion']]
        chart_data.add_series('Series 1', (float(ans[0]['clasificaciones'][0]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][1]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][2]['numero'])/float(total_clasificaciones)))

        x, y, cx, cy = Inches(6.69), Inches(4.7), Inches(3), Inches(2.5)

        chart = prs.slides[0].shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        #prs.save('visitas_stg/static/ppt/ppt-generados/Reporte_Dependencia_sisef_' + str(usuario.user.id) + '.pptx')
        #the_file = 'visitas_stg/static/ppt/ppt-generados/Reporte_Dependencia_sisef_' + str(usuario.user.id) + '.pptx'

        prs.save('/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Dependencia_sisef_' + str(usuario.user.id) + '.pptx')
        the_file = '/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Dependencia_sisef_' + str(usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response


def Predefinido_Region(request):

        #prs = Presentation('visitas_stg/static/ppt/Reporte_Region_sisef.pptx')
        prs = Presentation('/home/sisefenlin/visitas/static/ppt/Reporte_Region_sisef.pptx')
        usuario = request.user.userprofile

        print request.user.userprofile.rol

        if request.user.userprofile.rol == 'AD':
            dependencias = Dependencia.objects.all()
        else:
            dependencias = Dependencia.objects.filter(
                Q(id=request.user.userprofile.dependencia_id))

        ans = []

        date = datetime.now()
        json_fecha={}
        if date.day >= 10:
            json_fecha['dia'] = str(date.day)
        else:
            json_fecha['dia'] = "0" + str(date.day)
        if date.month >= 10:
            json_fecha['mes'] = str(date.month)
        else:
            json_fecha['mes'] = "0" + str(date.month)
        json_fecha['ano'] = str(date.year)

        region = Region.objects.get(id=request.GET.get('region_id'))
        medios = Medio.objects.values('id', 'nombre_medio')
        dependencias = Dependencia.objects.all()
        clasificaciones = Clasificacion.objects.values('id', 'nombre_clasificacion')

        map = {}
        map['region'] = region.to_serializable_dict()
        map['region']['estados'] = Estado.objects.filter(region_id=region.id).count()
        map['region']['distritos_electorales'] = DistritoElectoral.objects.filter(estado__region_id=region.id).count()
        map['region']['municipios'] = Municipio.objects.filter(estado__region_id=region.id).count()

        map['dependencias'] = []

        for dependencia in dependencias:
            dependencia_map = dependencia.to_serializable_dict()
            dependencia_map['funcionarios_federales'] = Cargo.objects.filter(dependencia_id=dependencia.id).count()
            dependencia_map['visitas'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(region_id=region.id)).count()
            dependencia_map['municipios'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=dependencia.id)).values(
                'municipio_id').distinct().count()
            dependencia_map['actividades'] = Actividad.objects.filter(
                Q(visita__dependencia_id=dependencia.id) & Q(visita__region_id=dependencia.id)).count()
            dependencia_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__region_id=region.id)).count()
            dependencia_map['capitalizaciones'] = Capitalizacion.objects.filter(
                Q(actividad__visita__region_id=region.id) & Q(
                    actividad__visita__dependencia_id=dependencia.id)).aggregate(Sum('cantidad'))
            dependencia_map['municipios'] = Visita.objects.filter(
                Q(region_id=region.id) & Q(dependencia=dependencia.id)).values('municipio_id').distinct().count()
            map['dependencias'].append(dependencia_map)

        map['dependencias'].sort(key=lambda e: e['capitalizaciones'])

        map['medios'] = []
        for medio in medios:
            medio_map = medio

            medio_map['tipos_capitalizacion'] = []
            tipos_capitalizacion = TipoCapitalizacion.objects.all()
            for tipo_capitalizacion in tipos_capitalizacion:
                tipo_map = tipo_capitalizacion.to_serializable_dict()
                tipo_map['numero'] = Capitalizacion.objects.filter(
                    Q(tipo_capitalizacion_id=tipo_capitalizacion.id) & Q(medio_id=medio['id']) & Q(
                        actividad__visita__region_id=region.id)).count()
                medio_map['tipos_capitalizacion'].append(tipo_map)
            map['medios'].append(medio_map)

        map['clasificaciones'] = []
        for clasificacion in clasificaciones:
            tipo_map = clasificacion
            tipo_map['numero'] = Actividad.objects.filter(
                Q(visita__region_id=region.id) & Q(clasificacion_id=clasificacion['id'])).count()
            map['clasificaciones'].append(tipo_map)

        ans.append(map)

        prs.slides[0].shapes[8].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[9].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[10].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.size = Pt(9)
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.size = Pt(9)


        prs.slides[0].shapes[8].text= json_fecha['dia'] + "/" + json_fecha['mes'] + "/" + json_fecha['ano']
        prs.slides[0].shapes[9].text= '{0:,}'.format(ans[0]['region']['distritos_electorales'])
        prs.slides[0].shapes[10].text= ans[0]['region']['numeroRegion']
        #prs.slides[0].shapes[7].text= '{0:,}'.format(ans[0]['estado']['distritos_electorales'])
        prs.slides[0].shapes[11].text= '{0:,}'.format(ans[0]['region']['municipios'])
        prs.slides[0].shapes[12].text= '{0:,}'.format(ans[0]['region']['estados'])


        table = prs.slides[0].shapes[0].table
        # write body cellstable.cell(1, 0)
        i=1

        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['funcionarios_federales'], reverse=True)
        total=0
        for dato in ans[0]['dependencias']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreDependencia'])
            table.cell(i, 1).text = str(dato['funcionarios_federales'])
            total+=dato['funcionarios_federales']
            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[1].table
        i=1
        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['visitas'], reverse=True)
        total=0
        for dato in ans[0]['dependencias']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreDependencia'])
            table.cell(i, 1).text = str(dato['visitas'])
            total += dato['visitas']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[2].table
        i=1
        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['actividades'], reverse=True)
        total=0
        for dato in ans[0]['dependencias']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreDependencia'])
            table.cell(i, 1).text = str(dato['actividades'])
            total+=dato['actividades']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[3].table
        i=1
        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['municipios'], reverse=True)
        total=0
        for dato in ans[0]['dependencias']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreDependencia'])
            table.cell(i, 1).text = str(dato['municipios'])
            total += dato['municipios']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[4].table
        i=1
        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['participantes_locales'], reverse=True)
        total=0
        for dato in ans[0]['dependencias']:
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i, 0).text = str(dato['nombreDependencia'])
            table.cell(i, 1).text = str(dato['participantes_locales'])
            total += dato['participantes_locales']

            if i==5: break
            i=i+1
        table.cell(6, 0).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 0).text = "TOTAL"
        table.cell(6, 1).text = str(total)

        table = prs.slides[0].shapes[5].table
        i=1
        total1=total2=total3=total4=total5=0
        for dato in ans[0]['medios']:
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,2).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,3).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,4).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,5).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,6).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,1).text = str(dato['tipos_capitalizacion'][0]['numero'])
            table.cell(i,2).text = str(dato['tipos_capitalizacion'][1]['numero'])
            table.cell(i,3).text = str(dato['tipos_capitalizacion'][2]['numero'])
            table.cell(i,4).text = str(dato['tipos_capitalizacion'][3]['numero'])
            table.cell(i,5).text = str(dato['tipos_capitalizacion'][4]['numero'])

            totalColumna=0
            for j in  range(0,5):
                totalColumna += dato['tipos_capitalizacion'][j]['numero']
            table.cell(i,6).text = str(totalColumna)

            total1=total1 + dato['tipos_capitalizacion'][0]['numero']
            total2=total2 + dato['tipos_capitalizacion'][1]['numero']
            total3=total3 + dato['tipos_capitalizacion'][2]['numero']
            total4=total4 + dato['tipos_capitalizacion'][3]['numero']
            total5=total5 + dato['tipos_capitalizacion'][4]['numero']
            i=i+1

        table.cell(6, 1).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 2).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 3).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 4).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 5).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 6).text_frame.paragraphs[0].font.size = Pt(8)
        table.cell(6, 1).text = str(total1)
        table.cell(6, 2).text = str(total2)
        table.cell(6, 3).text = str(total3)
        table.cell(6, 4).text = str(total4)
        table.cell(6, 5).text = str(total5)
        table.cell(6, 6).text = str(total1+total2+total3+total4+total5)

        mayor = map['dependencias']
        mayor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'], reverse=True)

        table = prs.slides[0].shapes[6].table
        i=1
        for dato in mayor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreDependencia'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        menor = map['dependencias']
        menor.sort(key=lambda x: x['capitalizaciones']['cantidad__sum'])

        table = prs.slides[0].shapes[7].table
        i=1
        for dato in menor:
            table.cell(i,0).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(i,1).text_frame.paragraphs[0].font.size = Pt(8)

            table.cell(i,0).text = str(dato['nombreDependencia'])
            table.cell(i,1).text = str(0)
            if str(dato['capitalizaciones']['cantidad__sum'])!='None':
                table.cell(i,1).text = str(dato['capitalizaciones']['cantidad__sum'])

            if i==3: break
            i=i+1

        total_clasificaciones=0
        for cantidad in ans[0]['clasificaciones']:
            total_clasificaciones = total_clasificaciones + cantidad['numero']

        #grafica pie
        chart_data = ChartData()
        chart_data.categories = [ans[0]['clasificaciones'][0]['nombre_clasificacion'], ans[0]['clasificaciones'][1]['nombre_clasificacion'], ans[0]['clasificaciones'][2]['nombre_clasificacion']]
        chart_data.add_series('Series 1', (float(ans[0]['clasificaciones'][0]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][1]['numero'])/float(total_clasificaciones), float(ans[0]['clasificaciones'][2]['numero'])/float(total_clasificaciones)))

        x, y, cx, cy = Inches(6.69), Inches(4.7), Inches(3), Inches(2.5)

        chart = prs.slides[0].shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        #prs.save('visitas_stg/static/ppt/ppt-generados/Reporte_Region_sisef_' + str(usuario.user.id) + '.pptx')
        #the_file = 'visitas_stg/static/ppt/ppt-generados/Reporte_Region_sisef_' + str(usuario.user.id) + '.pptx'

        prs.save('/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Region_sisef_' + str(usuario.user.id) + '.pptx')
        the_file = '/home/sisefenlin/visitas/static/ppt/ppt-generados/Reporte_Region_sisef_' + str(usuario.user.id) + '.pptx'

        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response


def ayuda(request):
    return render_to_response('admin/visitas_stg/ayuda/c_ayuda.html', locals(),
                              context_instance=RequestContext(request))
def videos(request):
    return render_to_response('admin/visitas_stg/videos/videos_lista.html', locals(),
                              context_instance=RequestContext(request))
def manualesPdf(request):
    return render_to_response('admin/visitas_stg/manuales/manuales_lista.html', locals(),
                              context_instance=RequestContext(request))
@login_required()
def ver_video(request):
    tituloVideo=""
    cualVideo=request.GET.get('cualVideo', None),
    print(str(cualVideo[0]))
    if str(cualVideo[0]) =='alta_VISITA.mp4':
        tituloVideo='Alta de una Visita',
    elif str(cualVideo[0]) =='modificar_VISITA.mp4':
        tituloVideo='Modificar una Visita',

    elif str(cualVideo[0]) =='consulta_FILTROS.mp4':
        tituloVideo='Consulta Mediante Filtros',
    elif str(cualVideo[0]) =='consulta_PRDETERMINADAS.mp4':
        tituloVideo='Consulta Predefinidos',
    elif str(cualVideo[0]) =='lista_VISITAS.mp4':
        tituloVideo='Consulta Predefinidos',

    elif str(cualVideo[0]) =='add_FUNCIONARIO.mp4':
        tituloVideo='Agregar un Funcionario',
    elif str(cualVideo[0]) =='seach_FUNCIONARIO.mp4':
        tituloVideo='Buscar un Funcionario',
    elif str(cualVideo[0]) =='modify_FUNCIONARIO.mp4':
        tituloVideo='Modificar un Funcionario',
    elif str(cualVideo[0]) =='delete_FUNCIONARIO.mp4':
        tituloVideo='Eliminar un Funcionario',

    elif str(cualVideo[0]) =='add_MEDIO.mp4':
        tituloVideo='Agregar un Medio',
    elif str(cualVideo[0]) =='seach_MEDIO.mp4':
        tituloVideo='Buscar un Medio',
    elif str(cualVideo[0]) =='modify_MEDIO.mp4':
        tituloVideo='Modificar un Medio',
    elif str(cualVideo[0]) =='delete_MEDIO.mp4':
        tituloVideo='Eliminar un Medio',

    elif str(cualVideo[0]) =='add_ACTIVIDAD.mp4':
        tituloVideo='Agregar un Tipo de Actividad',
    elif str(cualVideo[0]) =='seach_ACTIVIDAD.mp4':
        tituloVideo='Buscar un Tipo de Actividad',
    elif str(cualVideo[0]) =='modify_ACTIVIDAD.mp4':
        tituloVideo='Modificar un Tipo de Actividad',
    elif str(cualVideo[0]) =='delete_ACTIVIDAD.mp4':
        tituloVideo='Eliminar un Tipo de Actividad',

    elif str(cualVideo[0]) =='add_CAPITALIZACION.mp4':
        tituloVideo='Agregar un Tipo de Capitalizacion',
    elif str(cualVideo[0]) =='seach_CAPITALIZACION.mp4':
        tituloVideo='Buscar un Tipo de Capitalizacion',
    elif str(cualVideo[0]) =='modify_CAPITALIZACION.mp4':
        tituloVideo='Modificar un Tipo de Capitalizacion',
    elif str(cualVideo[0]) =='delete_CAPITALIZACION.mp4':
        tituloVideo='Eliminar un Tipo de Capitalizacion',

    template = loader.get_template('admin/visitas_stg/videos/videos_lista.html')
    context = RequestContext(request, {
        'cualVideo': cualVideo,
        'tituloVideo': tituloVideo,
    })
    return HttpResponse(template.render(context))

def redirect_admin(request):
    return redirect('admin/')