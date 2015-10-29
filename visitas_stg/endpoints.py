from django.http import HttpResponse
from oauth2_provider.views.generic import ProtectedResourceView
from BuscarVisitas import BuscarVisitas
from models import Estado, Municipio, TipoCapitalizacion, DistritoElectoral, Region, Cargo, Dependencia, \
    TipoActividad, Clasificacion, Medio, Visita, Capitalizacion, Actividad, ParticipanteLocal
import json
from oauth2_provider.models import AccessToken
from views import get_array_or_none,get_array_or_vacio
from oauth2_provider.models import AccessToken
from django.db.models import Sum, IntegerField, Q, Count
from django.db.models import F

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor

try:
    import cStringIO as StringIO
except ImportError:
    import StringIO
from xlsxwriter.workbook import Workbook
from django.core.servers.basehttp import FileWrapper
from django.http import StreamingHttpResponse


def get_usuario_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user.userprofile
    else:
        return None


class ReporteInicioEndpoint(ProtectedResourceView):
    def get(self, request):
        dependencia = get_usuario_for_token(request.GET.get('access_token')).dependencia
        if dependencia is None:
            visitas = Visita.objects.all()
            capitalizaciones = Capitalizacion.objects.all()
        else:
            visitas = Visita.objects.filter(dependencia=dependencia)
            capitalizaciones = Capitalizacion.objects.filter(actividad__visita__in=visitas)

        reporte = {}

        reporte['estados'] = []
        estados = Estado.objects.filter(~Q(id=33) & ~Q(id=34))
        for estado in estados:
            total_visitas = visitas.filter(entidad=estado).count()
            if total_visitas > 0:
                reporte_estado = {'estado': estado.to_serialzable_dict(),
                              'total_visitas': total_visitas }
                reporte['estados'].append(reporte_estado)

        reporte['medios'] = []
        for tipo_medio in Medio.objects.all():
            reporte_medio = {'medio': tipo_medio.to_serializable_dict(),
                             'total_apariciones': capitalizaciones.filter(medio=tipo_medio).aggregate(
                                 total=Sum('cantidad'))['total']}
            # 'total_apariciones': capitalizaciones.filter(medio=tipo_medio).count()}
            reporte['medios'].append(reporte_medio)

        return HttpResponse(json.dumps(reporte), 'application/json')


class ReporteEstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        ans = []

        estados = Estado.objects.values('id', 'nombreEstado')
        medios = Medio.objects.values('id', 'nombre_medio')
        dependencias = Dependencia.objects.values('id', 'nombreDependencia')

        for estado in estados:
            visitas = Visita.objects.filter(entidad_id=estado['id'])

            map = {}
            map['estado'] = estado
            map['estado']['distritos_electorales'] = DistritoElectoral.objects.filter(Q(estado_id=estado['id'])).count()
            map['estado']['distritos_electorales_visitados'] = visitas.values(
                'distrito_electoral_id').distinct().count()
            map['estado']['municipios'] = Municipio.objects.filter(estado_id=estado['id']).count()
            map['estado']['municipios_visitados'] = visitas.values('municipio_id').distinct().count()

            map['dependencias'] = []

            for dependencia in dependencias:
                dependencia_map = dependencia
                dependencia_map['funcionarios_federales'] = Cargo.objects.filter(
                    dependencia_id=dependencia['id']).count()
                dependencia_map['visitas'] = visitas.filter(dependencia_id=dependencia['id']).count()
                dependencia_map['actividades'] = Actividad.objects.filter(
                    Q(visita__dependencia_id=dependencia['id']) & Q(visita__entidad_id=estado['id'])).count()
                dependencia_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                    Q(actividad__visita__dependencia_id=dependencia['id']) & Q(
                        actividad__visita__entidad_id=estado['id'])).count()
                dependencia_map['capitalizaciones'] = Capitalizacion.objects.filter(
                    Q(actividad__visita__entidad_id=estado['id']) & Q(
                        actividad__visita__dependencia_id=dependencia['id'])).count()
                map['dependencias'].append(dependencia_map)

            map['medios'] = []
            for medio in medios:
                medio_map = medio

                medio_map['tipos_capitalizacion'] = []
                tipos_capitalizacion = TipoCapitalizacion.objects.values('id', 'nombre_tipo_capitalizacion')
                for tipo_capitalizacion in tipos_capitalizacion:
                    tipo_map = tipo_capitalizacion
                    tipo_map['numero'] = Capitalizacion.objects.filter(
                        Q(tipo_capitalizacion_id=tipo_capitalizacion['id']) & Q(medio_id=medio['id']) & Q(
                            actividad__visita__entidad_id=estado['id'])).count()
                    medio_map['tipos_capitalizacion'].append(tipo_map)
                map['medios'].append(medio_map)

            ans.append(map)

        return HttpResponse(json.dumps(ans), 'application/json')


class ReporteDependenciaEndpoint(ProtectedResourceView):
    def get(self, request):
        dependencia = Dependencia.objects.get(id=request.GET.get('dependencia_id'))
        visitas = Visita.objects.filter(dependencia_id=dependencia.id)
        clasificaciones = Clasificacion.objects.values('id', 'nombre_clasificacion')
        medios = Medio.objects.all()

        map = {}
        map['dependencia'] = dependencia.to_serializable_dict()
        map['dependencia']['estados_visitados'] = visitas.values('entidad_id').distinct().count()
        map['dependencia']['municipios_visitados'] = visitas.values('entidad_id').distinct().count()
        map['dependencia']['distritos_electorales_visitados'] = visitas.values(
            'distrito_electoral_id').distinct().count()

        map['estados'] = []
        estados = Estado.objects.all()
        for estado in estados:
            estado_map = estado.to_serialzable_dict()
            estado_map['total_visitas_funcionarios_federales'] = Visita.objects.filter(
                Q(cargo__dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            estado_map['total_visitas'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            estado_map['total_actividades'] = Actividad.objects.filter(
                Q(visita__entidad_id=estado.id) & Q(visita__dependencia_id=dependencia.id)).count()
            estado_map['municipios'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=dependencia.id)).values(
                'municipio_id').distinct().count()
            estado_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).count()
            estado_map['capitalizaciones'] = Capitalizacion.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).count()
            map['estados'].append(estado_map)
        map['estados'].sort(key=lambda e: e['capitalizaciones'])

        map['medios'] = []
        for medio in medios:
            medio_map = medio.to_serializable_dict()

            medio_map['tipos_capitalizacion'] = []
            tipos_capitalizacion = TipoCapitalizacion.objects.all()
            for tipo_capitalizacion in tipos_capitalizacion:
                tipo_map = tipo_capitalizacion.to_serializable_dict()
                tipo_map['numero'] = Capitalizacion.objects.filter(
                    Q(tipo_capitalizacion_id=tipo_capitalizacion.id) & Q(medio_id=medio.id) & Q(
                        actividad__visita__dependencia_id=dependencia.id)).count()
                medio_map['tipos_capitalizacion'].append(tipo_map)
            map['medios'].append(medio_map)

        map['clasificaciones'] = []
        for clasificacion in clasificaciones:
            tipo_map = clasificacion
            tipo_map['numero'] = Actividad.objects.filter(
                Q(visita__dependencia_id=dependencia.id) & Q(clasificacion_id=clasificacion['id'])).count()
            map['clasificaciones'].append(tipo_map)

        return HttpResponse(json.dumps(map), 'application/json')


class ReporteEstadoEndpoint(ProtectedResourceView):
    def get(self, request):
        estado = Estado.objects.get(id=request.GET.get('estado_id'))
        medios = Medio.objects.values('id', 'nombre_medio')
        dependencias = Dependencia.objects.all()
        clasificaciones = Clasificacion.objects.values('id', 'nombre_clasificacion')

        map = {}
        map['estado'] = estado.to_serialzable_dict()
        map['estado']['distritos_electorales'] = DistritoElectoral.objects.filter(estado_id=estado.id).count()
        map['estado']['municipios'] = Municipio.objects.filter(estado_id=estado.id).count()

        map['dependencias'] = []

        for dependencia in dependencias:
            dependencia_map = dependencia.to_serializable_dict()
            dependencia_map['funcionarios_federales'] = Cargo.objects.filter(dependencia_id=dependencia.id).count()
            dependencia_map['visitas'] = Visita.objects.filter(
                Q(dependencia_id=dependencia.id) & Q(entidad_id=estado.id)).count()
            dependencia_map['actividades'] = Actividad.objects.filter(
                Q(visita__dependencia_id=dependencia.id) & Q(visita__entidad_id=dependencia.id)).count()
            dependencia_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                Q(actividad__visita__dependencia_id=dependencia.id) & Q(
                    actividad__visita__entidad_id=estado.id)).count()
            dependencia_map['capitalizaciones'] = Capitalizacion.objects.filter(
                Q(actividad__visita__entidad_id=estado.id) & Q(
                    actividad__visita__dependencia_id=dependencia.id)).count()
            map['dependencias'].append(dependencia_map)

        map['medios'] = []
        for medio in medios:
            medio_map = medio

            medio_map['tipos_capitalizacion'] = []
            tipos_capitalizacion = TipoCapitalizacion.objects.all()
            for tipo_capitalizacion in tipos_capitalizacion:
                tipo_map = tipo_capitalizacion.to_serializable_dict()
                tipo_map['numero'] = Capitalizacion.objects.filter(
                    Q(tipo_capitalizacion_id=tipo_capitalizacion.id) & Q(medio_id=medio['id']) & Q(
                        actividad__visita__entidad_id=estado.id)).count()
                medio_map['tipos_capitalizacion'].append(tipo_map)
            map['medios'].append(medio_map)

        map['clasificaciones'] = []
        for clasificacion in clasificaciones:
            tipo_map = clasificacion
            tipo_map['numero'] = Actividad.objects.filter(
                Q(visita__estado_id=estado.id) & Q(clasificacion_id=clasificacion['id'])).count()
            map['clasificaciones'].append(tipo_map)

        # return HttpResponse(json.dumps(map), 'application/json')
        return HttpResponse(map.__str__(), 'application/json')


class ReporteDependenciasEndpoint(ProtectedResourceView):
    def get(self, request):
        ans = {}

        dependencias = Dependencia.objects.values('nombreDependencia', 'id')
        medios = Medio.objects.values('id', 'nombre_medio')
        tipos_capitalizacion = TipoCapitalizacion.objects.values('nombre_tipo_capitalizacion', 'id')

        ans['dependencias'] = []
        for dependencia in dependencias:
            visitas = Visita.objects.filter(dependencia_id=dependencia['id'])
            map = {}
            map['dependencia'] = dependencia
            map['dependencia']['estados_visitados'] = visitas.values('entidad_id').distinct().count()
            map['dependencia']['municipios_visitados'] = visitas.values('municipio_id').distinct().count()
            map['dependencia']['distritos_electorales_visitados'] = visitas.values(
                'distrito_electoral_id').distinct().count()

            estados = visitas.values('entidad_id', 'entidad__nombreEstado').distinct().annotate(
                total_visitas_funcionarios_federales=Count('id', distintct=True),
                total_visitas=Count('id', distinct=True),
                municipios=Count('municipio_id', distinct=True)
            )
            map['estados'] = []
            for estado in estados:
                estado_map = estado
                # estado_map['total_visitas_funcionarios_federales'] = Visita.objects.filter(
                # Q(cargo__dependencia_id=dependencia['id']) & Q(entidad_id=estado['id'])).count() # Maybe cannot be optimizer
                # estado_map['total_visitas'] = visitas.filter(Q(entidad_id=estado['entidad_id'])).count()
                estado_map['total_actividades'] = Actividad.objects.filter(
                    Q(visita__entidad_id=estado['entidad_id']) & Q(visita__dependencia_id=dependencia['id'])).count()
                # estado_map['municipios'] = Visita.objects.filter(
                # Q(dependencia_id=dependencia['id']) & Q(entidad_id=estado['id'])).values(
                # 'municipio_id').distinct().count()
                estado_map['participantes_locales'] = ParticipanteLocal.objects.filter(
                    Q(actividad__visita__dependencia_id=dependencia['id']) & Q(
                        actividad__visita__entidad_id=estado['entidad_id'])).count()
                estado_map['capitalizaciones'] = Capitalizacion.objects.filter(
                    Q(actividad__visita__dependencia_id=dependencia['id']) & Q(
                        actividad__visita__entidad_id=estado['entidad_id'])).count()
                map['estados'].append(estado_map)

            map['medios'] = []
            for medio in medios:
                medio_map = medio

                medio_map['tipos_capitalizacion'] = []
                # medio_map['tipos_capitalizacion'] = visitas.values(
                # 'actividad__capitalizacion__tipo_capitalizacion').distinct().annotate(
                # numero=Count('actividad__capitalizacion'))
                for tipo_capitalizacion in tipos_capitalizacion:
                    tipo_map = tipo_capitalizacion
                    tipo_map['numero'] = Capitalizacion.objects.filter(
                        Q(tipo_capitalizacion_id=tipo_capitalizacion['id']) & Q(medio_id=medio['id']) & Q(
                            actividad__visita__dependencia_id=dependencia['id'])).count()
                    medio_map['tipos_capitalizacion'].append(tipo_map)
                map['medios'].append(medio_map)

            ans['dependencias'].append(map)

        return HttpResponse(json.dumps(ans), 'application/json')


class IdUnicoEndpoint(ProtectedResourceView):
    def get(self, request):
        identificador_unico = request.GET.get('identificador_unico')
        ans = {'error': None, 'visita': None}
        if identificador_unico is not None:
            visita = Visita.objects.filter(identificador_unico=identificador_unico)
            if visita is not None and visita.count() > 0:
                visita = visita.first()
                dependencia_usuario = get_usuario_for_token(request.GET.get('access_token')).dependencia_id
                if dependencia_usuario is None or dependencia_usuario == visita.dependencia_id:
                    ans['visita'] = visita.first().to_serializable_dict()
                else:
                    ans['error'] = 'Privilegios insuficientes'
            else:
                ans['error'] = 'No se encontro la visita'
        else:
            ans['error'] = 'Debes ingresar un identificador unico'
        return HttpResponse(json.dumps(ans), 'application/json')


class RegionesEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda region: region.to_serializable_dict(), Region.objects.all())),
                            'application/json')


class EstadosForRegionesEndpoint(ProtectedResourceView):
    def get(self, request):
        region_ids = get_array_or_vacio(request.GET.get('regiones'))
        if region_ids is not None:
            estados = Estado.objects.filter(region_id__in=region_ids)
        else:
            estados = Estado.objects.all()

        return HttpResponse(json.dumps(map(lambda estado: estado.to_serialzable_dict(), estados)), 'application/json')


class MunicipiosForEstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        estado_ids = get_array_or_vacio(request.GET.get('estados'))
        if estado_ids is None or 33 in estado_ids or 34 in estado_ids:
            municipios = Municipio.objects.all()
        else:
            municipios = Municipio.objects.filter(estado_id__in=estado_ids)
        municipios.order_by('nombreMunicipio')

        the_list = []
        for municipio in municipios.values('id', 'nombreMunicipio'):
            the_list.append(municipio)

        return HttpResponse(json.dumps(the_list), 'application/json')


class DistritoElectoralForEstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        estado_ids = get_array_or_vacio(request.GET.get('estados'))

        if estado_ids is not None:
            distritos = DistritoElectoral.objects.filter(estado_id__in=estado_ids)
        else:
            distritos = DistritoElectoral.objects.all()

        return HttpResponse(json.dumps(map(lambda distrito: distrito.to_serializable_dict(), distritos)),
                            'application/json')


class CargosForDependenciasEndpoint(ProtectedResourceView):
    def get(self, request):
        dependencia_ids = get_array_or_vacio(request.GET.get('dependencias',None))

        if dependencia_ids is not None and len(dependencia_ids) > 0:
            cargos = Cargo.objects.filter(dependencia_id__in=dependencia_ids)
        else:
            cargos = Cargo.objects.all()

        return HttpResponse(json.dumps(map(lambda cargo: cargo.to_serializable_dict(), cargos)), 'application/json')


class CargosForNombreEndpoint(ProtectedResourceView):
    def get(self, request):
        cargo_ids = get_array_or_vacio(request.GET.get('id_cargo'))

        if cargo_ids is not None and len(cargo_ids) > 0:
            cargos = Cargo.objects.filter(id__in=cargo_ids)
        else:
            cargos = Cargo.objects.all()

        return HttpResponse(json.dumps(map(lambda cargo: cargo.to_serializable_dict(), cargos)), 'application/json')


class CargosForCargosEndpoint(ProtectedResourceView):
    def get(self, request):
        cargos_string = request.GET.get('cargos', None)

        if cargos_string is None:
            cargos_nombres = None
        else:
            cargos_nombres = cargos_string.split(',')

        if cargos_nombres is not None and len(cargos_nombres) > 0:
            cargos = Cargo.objects.all()
        else:
            cargos = Cargo.objects.filter(nombre_cargo__in=cargos_nombres)

        return HttpResponse(json.dumps(map(lambda cargo: cargo.to_serializable_dict(), cargos)), 'application/json')


class IdUnicoEndpoint(ProtectedResourceView):
    def get(self, request):
        identificador_unico = request.GET.get('identificador_unico')
        ans = {'error': None, 'visita': None}
        if identificador_unico is not None:
            visita = Visita.objects.filter(identificador_unico=identificador_unico)
            if visita is not None and visita.count() > 0:
                visita = visita.first()
                dependencia_usuario = get_usuario_for_token(request.GET.get('access_token')).dependencia_id
                if dependencia_usuario is None or dependencia_usuario == visita.dependencia_id:
                    ans['visita'] = visita.to_serializable_dict()
                else:
                    ans['error'] = 'Privilegios insuficientes'
            else:
                ans['error'] = 'No se encontro la visita'
        else:
            ans['error'] = 'Debes ingresar un identificador unico'
        return HttpResponse(json.dumps(ans), 'application/json')


class BuscarVisitasEndpoint(ProtectedResourceView):
    def get(self, request):
        buscador = BuscarVisitas(ids_dependencia=get_array_or_none(request.GET.get('dependencia')),
                                 rango_fecha_inicio=request.GET.get('fechaInicio', None),
                                 rango_fecha_fin=request.GET.get('fechaFin', None),
                                 descripcion=request.GET.get('descripcion', None),
                                 problematica=request.GET.get('problematica', None),
                                 nombre_medio=request.GET.get('nombreMedio', None),
                                 ids_region=get_array_or_none(request.GET.get('region')),
                                 ids_entidad=get_array_or_none(request.GET.get('estado')),
                                 ids_municipio=get_array_or_none(request.GET.get('municipio')),
                                 ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                                 ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                                 ids_partido=get_array_or_none(request.GET.get('partido')),
                                 identificador_unico=request.GET.get('identificador_unico', None),
                                 ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                                 ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                                 ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                                 ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                                 limite_min=request.GET.get('limiteMin', 0),
                                 limite_max=request.GET.get('limiteMax', 100),
                                 )
        ans = buscador.buscar()

        json_ans = {}
        json_ans['visitas'] = []
        for visita in ans['visitas']:
            json_ans['visitas'].append(visita.to_serializable_dict())

        json_ans['reporte_general'] = {'visitas_totales': ans['reporte_general']['visitas_totales']}

        json_ans['reporte_estado'] = []
        for estado in ans['reporte_estado']:
            json_map = {'estado': estado['entidad__nombreEstado'],
                        'numero_visitas': estado['numero_visitas'],
                        'numero_apariciones': estado['numero_apariciones']}
            json_ans['reporte_estado'].append(json_map)

        json_ans['reporte_dependencia'] = []
        for dependencia in ans['reporte_dependencia']:
            json_map = {'dependencia': dependencia['dependencia__nombreDependencia'],
                        'numero_visitas': dependencia['numero_visitas'],
                        'numero_apariciones': dependencia['numero_apariciones']}
            json_ans['reporte_dependencia'].append(json_map)

        return HttpResponse(json.dumps(json_ans), 'application/json')


class DependenciasEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.all())),
            'application/json')


class TipoActividadEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda tipo: tipo.to_serializable_dict(), TipoActividad.objects.all())),
                            "application/json")


class ClasificacionEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda clasificacion: clasificacion.to_serializable_dict(), Clasificacion.objects.all())),
            'application/json')


class MediosEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda medio: medio.to_serializable_dict(), Medio.objects.all())),
                            'application/json')


class TipoCapitalizacionEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda tipoCapitalizacion: tipoCapitalizacion.to_serializable_dict(),
                           TipoCapitalizacion.objects.all())), 'application/json')


class PptxEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user

        buscador = BuscarVisitas(ids_dependencia=get_array_or_none(request.GET.get('dependencia')),
                                 rango_fecha_inicio=request.GET.get('fechaInicio', None),
                                 rango_fecha_fin=request.GET.get('fechaFin', None),
                                 descripcion=request.GET.get('descripcion', None),
                                 problematica=request.GET.get('problematica', None),
                                 nombre_medio=request.GET.get('nombreMedio', None),
                                 ids_region=get_array_or_none(request.GET.get('region')),
                                 ids_entidad=get_array_or_none(request.GET.get('estado')),
                                 ids_municipio=get_array_or_none(request.GET.get('municipio')),
                                 ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                                 ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                                 ids_partido=get_array_or_none(request.GET.get('partido')),
                                 identificador_unico=request.GET.get('identificador_unico', None),
                                 ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                                 ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                                 ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                                 ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                                 limite_min=request.GET.get('limiteMin', 0),
                                 limite_max=request.GET.get('limiteMax', 100),
                                 )
        resultados = buscador.buscar()

        json_map = {}

        json_map['visitas'] = []
        for visita in resultados['visitas'].values('id', 'identificador_unico', 'entidad__nombreEstado',
                                                   'actividad__descripcion'):
            json_map['visitas'].append(visita)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Resultados'

        renglones = resultados['reporte_general']['visitas_totales'] + 1
        if renglones < 22:
            rows = renglones
        else:
            rows = 22
        cols = 3
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(4.0)

        # write column headings
        table.cell(0, 0).text = 'Identificador'
        table.cell(0, 1).text = 'Actividad'
        table.cell(0, 2).text = 'Estado'

        # write body cells
        indice = 1
        for obra in json_map['visitas']:

            if indice == 22:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 22
                cols = 3
                left = Inches(0.921)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table
                # set column widths
                table.columns[0].width = Inches(2.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(4.0)

            # write column headings
            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table.cell(0, 0).text = 'Identificador'
            table.cell(0, 1).text = 'Actividad'
            table.cell(0, 2).text = 'Estado'

            # write body cells
            table.cell(indice, 0).text = obra['identificador_unico']
            table.cell(indice, 1).text = obra['actividad__descripcion']
            table.cell(indice, 2).text = obra['entidad__nombreEstado']
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="resultado_visitas.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class PptxReporteEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user
        tipoReporte = request.GET.get("tipoReporte", None)
        buscador = BuscarVisitas(ids_dependencia=get_array_or_none(request.GET.get('dependencia')),
                                 rango_fecha_inicio=request.GET.get('fechaInicio', None),
                                 rango_fecha_fin=request.GET.get('fechaFin', None),
                                 descripcion=request.GET.get('descripcion', None),
                                 problematica=request.GET.get('problematica', None),
                                 nombre_medio=request.GET.get('nombreMedio', None),
                                 ids_region=get_array_or_none(request.GET.get('region')),
                                 ids_entidad=get_array_or_none(request.GET.get('estado')),
                                 ids_municipio=get_array_or_none(request.GET.get('municipio')),
                                 ids_cargo_ejecuta=get_array_or_none(request.GET.get('cargoEjecuta')),
                                 ids_distrito_electoral=get_array_or_none(request.GET.get('distritoElectoral')),
                                 ids_partido=get_array_or_none(request.GET.get('partido')),
                                 identificador_unico=request.GET.get('identificador_unico', None),
                                 ids_tipo_actividad=get_array_or_none(request.GET.get('tipoActividad')),
                                 ids_tipo_medio=get_array_or_none(request.GET.get('tipoMedio')),
                                 ids_tipo_capitalizacion=get_array_or_none(request.GET.get('tipoCapitalizacion')),
                                 ids_clasificacion=get_array_or_none(request.GET.get('tipoClasificacion')),
                                 limite_min=request.GET.get('limiteMin', 0),
                                 limite_max=request.GET.get('limiteMax', 100),
                                 )
        resultados = buscador.buscar()

        json_map = {}

        json_map['reporte_dependencia'] = []
        for reporte in resultados['reporte_dependencia']:
            map = {}
            map['dependencia'] = Dependencia.objects.get(
                nombreDependencia=reporte['dependencia__nombreDependencia']).to_serializable_dict()
            map['numero_visitas'] = reporte['numero_visitas']
            if resultados['reporte_general']['visitas_totales'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(resultados['reporte_general']['visitas_totales'])
            json_map['reporte_dependencia'].append(map)

        json_map['reporte_estado'] = []
        for reporte_estado in resultados['reporte_estado']:
            map = {}
            if resultados['reporte_general']['visitas_totales'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(resultados['reporte_general']['visitas_totales'])
            map['estado'] = Estado.objects.get(
                nombreEstado=reporte_estado['entidad__nombreEstado'])
            map['numero_visitas'] = reporte_estado['numero_visitas']

            json_map['reporte_estado'].append(map)

        json_map['reporte_general'] = []
        map['visitas_totales'] = resultados['reporte_general']['visitas_totales']
        json_map['reporte_general'].append(map)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Reporte'

        if tipoReporte == 'Dependencia':


            rows = 17
            cols = 3
            left = Inches(0.921)
            top = Inches(1.2)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(2.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # write column headings
            table.cell(0, 0).text = 'Origen'
            table.cell(0, 1).text = 'No. de Visitas '
            table.cell(0, 2).text = 'Monto'

            # write body cells
            i = 1
            for obra in json_map['reporte_dependencia']:
                for x in range(0, 3):
                    cell = table.rows[i].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                table.cell(i, 0).text = obra['dependencia']['nombreDependencia']
                table.cell(i, 1).text = str(obra['numero_visitas'])
                table.cell(i, 2).text = str(0)
                i += 1

        if tipoReporte == 'Estado':
            rows = 35
            cols = 3
            left = Inches(0.921)
            top = Inches(1.2)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(2.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # write column headings
            table.cell(0, 0).text = 'Origen'
            table.cell(0, 1).text = 'No. de visitas '
            table.cell(0, 2).text = 'Monto'

            # write body cells
            i = 1
            for obra in json_map['reporte_estado']:
                for x in range(0, 3):
                    cell = table.rows[i].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
                table.cell(i, 0).text = obra['estado'].nombreEstado
                table.cell(i, 1).text = str(obra['numero_visitas'])
                table.cell(i, 2).text = str(0)
                i += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="resultado_visitas.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response
